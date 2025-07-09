#!/usr/bin/env python3
"""
CSV-based PPT Generator
Reads data from CSV files and generates PowerPoint presentations with charts and insights
"""

import os
import re
import json
from datetime import datetime
from typing import Dict, Any, List, Optional

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')  # Use non-GUI backend for web environments
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.ticker import FixedLocator

# Excel support
import openpyxl
from openpyxl import load_workbook
import xlrd

import openai
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

class CSVPPTGenerator:
    def __init__(self):
        """Initialize the CSV PPT Generator with OpenAI client"""
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key or api_key == 'your_openai_api_key_here':
            raise ValueError("Please set your OpenAI API key in the .env file")
        self.openai_client = openai.OpenAI(api_key=api_key)
        sns.set_palette("husl")
        self.data_analysis = {}
        self.charts_created = []

    def detect_file_type(self, file_path: str) -> str:
        """Detect if file is CSV or Excel"""
        file_extension = os.path.splitext(file_path)[1].lower()
        if file_extension in ['.xlsx', '.xls']:
            return 'excel'
        elif file_extension == '.csv':
            return 'csv'
        else:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: .csv, .xlsx, .xls")
    
    def load_excel_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about Excel file (sheets, named ranges)"""
        try:
            # Load workbook to get sheet information
            wb = load_workbook(file_path, read_only=True)
            
            sheet_info = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                # Get sheet dimensions
                max_row = ws.max_row
                max_col = ws.max_column
                
                # Check if sheet has data
                has_data = max_row > 1 or max_col > 1
                
                sheet_info[sheet_name] = {
                    'max_row': max_row,
                    'max_col': max_col,
                    'has_data': has_data,
                    'estimated_records': max_row - 1 if has_data else 0  # Subtract header row
                }
            
            # Get named ranges
            named_ranges = []
            try:
                for name, definition in wb.defined_names.definedName:
                    if not definition.is_reserved:
                        named_ranges.append({
                            'name': name,
                            'range': str(definition.value),
                            'sheet': definition.destinations[0][0] if definition.destinations else 'Unknown'
                        })
            except Exception as e:
                print(f"Warning: Could not read named ranges: {e}")
            
            wb.close()
            
            return {
                'file_path': file_path,
                'sheets': sheet_info,
                'named_ranges': named_ranges,
                'total_sheets': len(sheet_info),
                'sheets_with_data': len([s for s in sheet_info.values() if s['has_data']])
            }
            
        except Exception as e:
            raise ValueError(f"Error reading Excel file information: {e}")
    
    def load_excel_sheet(self, file_path: str, sheet_name: str = None, named_range: str = None) -> pd.DataFrame:
        """Load specific sheet or named range from Excel file"""
        try:
            if named_range:
                # Load specific named range
                df = pd.read_excel(file_path, sheet_name=sheet_name, usecols=named_range)
                print(f"📊 Loaded named range '{named_range}' from sheet '{sheet_name}'")
            elif sheet_name:
                # Load specific sheet
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                print(f"📊 Loaded sheet '{sheet_name}' ({len(df)} rows, {len(df.columns)} columns)")
            else:
                # Load first sheet by default
                df = pd.read_excel(file_path, sheet_name=0)
                print(f"📊 Loaded first sheet ({len(df)} rows, {len(df.columns)} columns)")
            
            # Basic validation
            if df.empty:
                raise ValueError("The selected sheet/range is empty")
            
            # Handle common Excel issues
            # Remove completely empty rows and columns
            df = df.dropna(how='all').dropna(axis=1, how='all')
            
            if df.empty:
                raise ValueError("No data found after removing empty rows/columns")
            
            print(f"✅ Data loaded successfully: {len(df)} rows × {len(df.columns)} columns")
            return df
            
        except Exception as e:
            raise ValueError(f"Error loading Excel data: {e}")
    
    def choose_best_sheet(self, excel_info: Dict[str, Any]) -> str:
        """Automatically choose the best sheet to analyze"""
        sheets = excel_info['sheets']
        
        # Priority 1: Sheet with most data
        best_sheet = max(sheets.keys(), key=lambda s: sheets[s]['estimated_records'])
        
        # Priority 2: Avoid sheets that might be summaries or metadata
        avoid_names = ['summary', 'metadata', 'info', 'readme', 'instructions']
        data_sheets = [name for name in sheets.keys() 
                      if not any(avoid in name.lower() for avoid in avoid_names) 
                      and sheets[name]['has_data']]
        
        if data_sheets:
            best_sheet = max(data_sheets, key=lambda s: sheets[s]['estimated_records'])
        
        print(f"🎯 Auto-selected sheet: '{best_sheet}' ({sheets[best_sheet]['estimated_records']} estimated records)")
        return best_sheet
    
    def load_and_analyze_data(self, file_path: str, sheet_name: str = None, named_range: str = None) -> Dict[str, Any]:
        """Load and analyze data from CSV or Excel file"""
        file_type = self.detect_file_type(file_path)
        
        if file_type == 'excel':
            return self.load_and_analyze_excel(file_path, sheet_name, named_range)
        else:
            return self.load_and_analyze_csv(file_path)
    
    def load_and_analyze_excel(self, file_path: str, sheet_name: str = None, named_range: str = None) -> Dict[str, Any]:
        """Load Excel file and perform comprehensive analysis"""
        try:
            # Get Excel file information
            excel_info = self.load_excel_info(file_path)
            print(f"📁 Excel file info: {excel_info['total_sheets']} sheets, {excel_info['sheets_with_data']} with data")
            
            # If no sheet specified, choose the best one
            if not sheet_name:
                sheet_name = self.choose_best_sheet(excel_info)
            
            # Load the data
            df = self.load_excel_sheet(file_path, sheet_name, named_range)
            
            # Store Excel-specific metadata
            excel_metadata = {
                'source_type': 'excel',
                'source_sheet': sheet_name,
                'source_named_range': named_range,
                'excel_info': excel_info
            }
            
            # Perform standard analysis
            analysis = self._perform_data_analysis(df, file_path, excel_metadata)
            return analysis
            
        except Exception as e:
            raise ValueError(f"Error loading Excel file: {e}")
    def load_and_analyze_csv(self, csv_file_path: str) -> Dict[str, Any]:
        """Load CSV file and perform comprehensive analysis with data cleaning"""
        try:
            # Load and clean CSV data
            df = self._load_csv_with_cleaning(csv_file_path)
            
            # CSV-specific metadata
            csv_metadata = {
                'source_type': 'csv',
                'source_sheet': None,
                'source_named_range': None,
                'excel_info': None
            }
            
            # Perform standard analysis
            analysis = self._perform_data_analysis(df, csv_file_path, csv_metadata)
            return analysis
            
        except Exception as e:
            raise ValueError(f"Error loading CSV file: {e}")
    def _clean_data_for_perfect_ppt(self, df: pd.DataFrame) -> pd.DataFrame:
        """Advanced data cleaning for perfect PPT generation"""
        print("\n🔧 Advanced Data Cleaning for Perfect PPT...")
        
        # 1. Handle text data quality
        print("📝 Cleaning text data...")
        string_columns = df.select_dtypes(include=['object']).columns
        for col in string_columns:
            # Remove extra whitespace and normalize case
            df[col] = df[col].astype(str).str.strip().str.replace(r'\s+', ' ', regex=True)
            # Replace 'nan' strings with actual NaN
            df[col] = df[col].replace(['nan', 'NaN', 'NULL', 'null', 'None'], pd.NA)
            # Remove rows with placeholder values
            df = df[~df[col].str.contains(r'^[0-9a-f]{32}$', regex=True, na=False)]  # Remove hash-like IDs
        
        # 2. Handle numeric data quality
        print("🔢 Cleaning numeric data...")
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            # Remove infinite values
            df = df[~np.isinf(df[col])]
            # Cap extreme values (beyond 3 standard deviations)
            mean_val = df[col].mean()
            std_val = df[col].std()
            if std_val > 0:
                lower_bound = mean_val - 3 * std_val
                upper_bound = mean_val + 3 * std_val
                extreme_count = len(df[(df[col] < lower_bound) | (df[col] > upper_bound)])
                if extreme_count > 0:
                    print(f"  ⚠️ Capping {extreme_count} extreme values in {col}")
                    df[col] = df[col].clip(lower=lower_bound, upper=upper_bound)
        
        # 3. Handle date/time data quality
        print("📅 Cleaning date/time data...")
        datetime_cols = df.select_dtypes(include=['datetime']).columns
        for col in datetime_cols:
            # Remove future dates if they seem unrealistic
            future_threshold = pd.Timestamp.now() + pd.DateOffset(years=5)
            future_dates = df[col] > future_threshold
            if future_dates.any():
                print(f"  ⚠️ Removing {future_dates.sum()} unrealistic future dates in {col}")
                df = df[~future_dates]
        
        # 4. Ensure data consistency
        print("🔄 Ensuring data consistency...")
        # Remove rows where key business rules are violated
        # Example: if you have Opening Balance + Allocate Budget != Closing Balance
        if all(col in df.columns for col in ['Opening Balance', 'Allocate Budget', 'Closing Balance']):
            balance_check = abs(df['Opening Balance'] + df['Allocate Budget'] - df['Closing Balance']) > 0.01
            inconsistent_rows = balance_check.sum()
            if inconsistent_rows > 0:
                print(f"  ⚠️ Found {inconsistent_rows} rows with balance inconsistencies")
                # You can choose to fix or remove these rows
                # For now, we'll keep them but flag them
        
        # 5. Final quality checks
        print("✅ Final quality validation...")
        # Ensure minimum data quality standards
        min_rows = 5  # Minimum rows for meaningful analysis
        if len(df) < min_rows:
            raise ValueError(f"Insufficient data after cleaning. Need at least {min_rows} rows, got {len(df)}")
        
        # Ensure we have at least some numeric data for charts
        if len(df.select_dtypes(include=[np.number]).columns) == 0:
            print("  ⚠️ No numeric columns found - limited chart options available")
        
        return df
    
    def _load_csv_with_cleaning(self, csv_file_path: str) -> pd.DataFrame:
        """Load CSV data with comprehensive cleaning for analysis"""
        try:
            print(f"🧹 Loading and cleaning CSV file: {csv_file_path}")
            
            # Try different encodings if utf-8 fails
            encodings = ['utf-8', 'utf-8-sig', 'latin1', 'iso-8859-1']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(csv_file_path, encoding=encoding)
                    print(f"✅ Successfully loaded with {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise ValueError("Could not load CSV file with any supported encoding")
            
            print(f"📊 Original data shape: {df.shape}")
            
            # Clean column names - remove leading/trailing whitespace and BOM
            df.columns = df.columns.str.strip().str.replace('\ufeff', '')
            
            # Clean string data - remove leading/trailing whitespace
            string_columns = df.select_dtypes(include=['object']).columns
            for col in string_columns:
                df[col] = df[col].astype(str).str.strip()
                # Replace empty strings with NaN
                df[col] = df[col].replace('', pd.NA)
            
            # Handle missing values more intelligently
            print(f"🔍 Missing values before cleaning: {df.isnull().sum().sum()}")
            
            # Drop rows where all values are missing
            df = df.dropna(how='all')
            
            # Drop columns where all values are missing
            df = df.dropna(axis=1, how='all')
            
            # For critical columns, fill missing values or drop rows
            # This can be customized based on your data requirements
            initial_rows = len(df)
            
            # Remove rows with missing values in key columns (if any)
            df = df.dropna(subset=df.columns[:3])  # Keep rows with data in first 3 columns
            
            print(f"📉 Rows removed due to missing key data: {initial_rows - len(df)}")
            
            # Remove duplicates
            initial_rows = len(df)
            df = df.drop_duplicates()
            duplicates_removed = initial_rows - len(df)
            print(f"🔄 Duplicate rows removed: {duplicates_removed}")
            
            # Convert date columns if detected
            for col in df.columns:
                if 'date' in col.lower():
                    try:
                        df[col] = pd.to_datetime(df[col], errors='coerce')
                        print(f"📅 Converted {col} to datetime")
                    except:
                        pass
            
            # Handle numeric columns
            for col in df.columns:
                if df[col].dtype == 'object':
                    # Try to convert to numeric if it looks like numbers
                    numeric_series = pd.to_numeric(df[col], errors='coerce')
                    if not numeric_series.isna().all():
                        non_na_count = numeric_series.notna().sum()
                        total_count = len(df)
                        if non_na_count / total_count > 0.5:  # If >50% can be converted
                            df[col] = numeric_series
                            print(f"🔢 Converted {col} to numeric ({non_na_count}/{total_count} values)")
            
            # Identify and handle outliers using IQR method
            print("\n🔍 Detecting outliers...")
            outlier_summary = {}
            for col in df.select_dtypes(include=[np.number]).columns:
                Q1 = df[col].quantile(0.25)
                Q3 = df[col].quantile(0.75)
                IQR = Q3 - Q1
                
                if IQR > 0:  # Only process if there's variation
                    # Count outliers but don't remove them yet
                    outlier_condition = (df[col] < (Q1 - 1.5 * IQR)) | (df[col] > (Q3 + 1.5 * IQR))
                    outlier_count = outlier_condition.sum()
                    outlier_summary[col] = outlier_count
                    
                    if outlier_count > 0:
                        print(f"  📊 {outlier_count} outliers detected in {col}")
                        # Only remove outliers if they're excessive (>10% of data)
                        if outlier_count / len(df) > 0.1:
                            print(f"    🗑️  Removing {outlier_count} outliers (>10% of data)")
                            df = df[~outlier_condition]
                        else:
                            print(f"    ✅ Keeping outliers (≤10% of data)")
            
            # Apply advanced cleaning for perfect PPT
            df = self._clean_data_for_perfect_ppt(df)
            
            print(f"\n✅ Final cleaned data shape: {df.shape}")
            
            if df.empty:
                raise ValueError("CSV file is empty after cleaning")
            
            # Show comprehensive data quality summary
            print(f"\n📈 Comprehensive Data Quality Summary:")
            print(f"  - Total rows: {len(df):,}")
            print(f"  - Total columns: {len(df.columns)}")
            print(f"  - Missing values: {df.isnull().sum().sum()}")
            print(f"  - Numeric columns: {len(df.select_dtypes(include=[np.number]).columns)}")
            print(f"  - Text columns: {len(df.select_dtypes(include=['object']).columns)}")
            print(f"  - DateTime columns: {len(df.select_dtypes(include=['datetime']).columns)}")
            
            # Data completeness percentage
            total_cells = len(df) * len(df.columns)
            missing_cells = df.isnull().sum().sum()
            completeness = ((total_cells - missing_cells) / total_cells * 100) if total_cells > 0 else 0
            print(f"  - Data completeness: {completeness:.1f}%")
            
            # Show outlier summary
            if outlier_summary:
                total_outliers = sum(outlier_summary.values())
                print(f"  - Total outliers detected: {total_outliers}")
            
            return df
            
        except Exception as e:
            raise ValueError(f"Error cleaning CSV file: {e}")
    
    def _perform_data_analysis(self, df: pd.DataFrame, file_path: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Perform comprehensive analysis on loaded data"""
        try:
            # Basic information
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
            datetime_cols = df.select_dtypes(include=['datetime']).columns.tolist()
            
            # Advanced analysis
            missing_info = df.isnull().sum()
            missing_percentage = (missing_info / len(df) * 100).round(2)
            
            # Data quality assessment
            data_quality = {
                "total_missing": missing_info.sum(),
                "missing_percentage_by_column": missing_percentage.to_dict(),
                "columns_with_missing": missing_info[missing_info > 0].index.tolist(),
                "complete_rows": len(df.dropna()),
                "duplicate_rows": df.duplicated().sum()
            }
            
            # Statistical insights for numeric columns
            numeric_insights = {}
            if numeric_cols:
                for col in numeric_cols:
                    col_data = df[col].dropna()
                    if len(col_data) > 0:
                        numeric_insights[col] = {
                            "mean": col_data.mean(),
                            "median": col_data.median(),
                            "std": col_data.std(),
                            "min": col_data.min(),
                            "max": col_data.max(),
                            "range": col_data.max() - col_data.min(),
                            "skewness": col_data.skew(),
                            "outliers_count": self._count_outliers(col_data)
                        }
            
            # Categorical insights
            categorical_insights = {}
            if categorical_cols:
                for col in categorical_cols:
                    col_data = df[col].dropna()
                    if len(col_data) > 0:
                        value_counts = col_data.value_counts()
                        categorical_insights[col] = {
                            "unique_count": col_data.nunique(),
                            "most_frequent": value_counts.index[0] if len(value_counts) > 0 else None,
                            "most_frequent_count": value_counts.iloc[0] if len(value_counts) > 0 else 0,
                            "distribution": value_counts.head(5).to_dict(),
                            "concentration": (value_counts.iloc[0] / len(col_data) * 100).round(2) if len(value_counts) > 0 else 0
                        }
            
            # Correlation analysis (enhanced)
            correlations = {}
            strong_correlations = []
            if len(numeric_cols) >= 2:
                corr_matrix = df[numeric_cols].corr()
                correlations = corr_matrix.to_dict()
                
                # Find strong correlations (>0.7 or <-0.7)
                for i in range(len(numeric_cols)):
                    for j in range(i+1, len(numeric_cols)):
                        corr_val = corr_matrix.iloc[i, j]
                        if abs(corr_val) > 0.7:
                            strong_correlations.append({
                                "var1": numeric_cols[i],
                                "var2": numeric_cols[j],
                                "correlation": round(corr_val, 3),
                                "strength": "strong positive" if corr_val > 0.7 else "strong negative"
                            })
            
            # Data patterns and trends
            patterns = self._identify_data_patterns(df, numeric_cols, categorical_cols)
            
            # Create analysis with source metadata
            analysis = {
                "file_name": os.path.basename(file_path),
                "shape": df.shape,
                "columns": df.columns.tolist(),
                "dtypes": df.dtypes.astype(str).to_dict(),
                "missing_values": missing_info.to_dict(),
                "numeric_columns": numeric_cols,
                "categorical_columns": categorical_cols,
                "datetime_columns": datetime_cols,
                "data_quality": data_quality,
                "numeric_insights": numeric_insights,
                "categorical_insights": categorical_insights,
                "summary_stats": df.describe().to_dict() if numeric_cols else {},
                "correlations": correlations,
                "strong_correlations": strong_correlations,
                "data_patterns": patterns,
                "sample_data": df.head(3).to_dict('records'),
                "source_metadata": metadata  # Include Excel/CSV metadata
            }
            
            self.df = df
            self.data_analysis = analysis
            return analysis
            
        except Exception as e:
            raise ValueError(f"Error performing data analysis: {e}")
    
    def _count_outliers(self, data: pd.Series) -> int:
        """Count outliers using IQR method"""
        Q1 = data.quantile(0.25)
        Q3 = data.quantile(0.75)
        IQR = Q3 - Q1
        lower_bound = Q1 - 1.5 * IQR
        upper_bound = Q3 + 1.5 * IQR
        outliers = data[(data < lower_bound) | (data > upper_bound)]
        return len(outliers)
    
    def _identify_data_patterns(self, df: pd.DataFrame, numeric_cols: List[str], categorical_cols: List[str]) -> Dict[str, Any]:
        """Identify interesting patterns in the data"""
        patterns = {
            "data_skewness": {},
            "potential_time_series": [],
            "high_cardinality_categories": [],
            "potential_ids": [],
            "suggested_groupings": []
        }
        
        # Check for skewness in numeric data
        for col in numeric_cols:
            col_data = df[col].dropna()
            if len(col_data) > 3:
                skew_val = col_data.skew()
                if abs(skew_val) > 1:
                    patterns["data_skewness"][col] = {
                        "skewness": round(skew_val, 3),
                        "interpretation": "highly skewed" if abs(skew_val) > 2 else "moderately skewed"
                    }
        
        # Check for potential time series columns
        for col in df.columns:
            col_lower = col.lower()
            if any(keyword in col_lower for keyword in ['date', 'time', 'year', 'month', 'day']):
                patterns["potential_time_series"].append(col)
        
        # Check for high cardinality categorical columns (potential IDs)
        for col in categorical_cols:
            unique_ratio = df[col].nunique() / len(df)
            if unique_ratio > 0.8:
                patterns["potential_ids"].append(col)
            elif df[col].nunique() > 20:
                patterns["high_cardinality_categories"].append(col)
        
        # Suggest meaningful groupings
        if categorical_cols and numeric_cols:
            for cat_col in categorical_cols[:2]:  # Limit to first 2 categorical columns
                for num_col in numeric_cols[:2]:  # Limit to first 2 numeric columns
                    patterns["suggested_groupings"].append({
                        "group_by": cat_col,
                        "analyze": num_col,
                        "chart_type": "bar"
                    })
        
        return patterns
    
    def _build_comprehensive_data_summary(self, analysis: Dict[str, Any]) -> str:
        """Build comprehensive data summary for AI to provide better context"""
        summary_parts = []
        
        # Basic dataset information
        summary_parts.append(f"📊 DATASET OVERVIEW:")
        summary_parts.append(f"• File: {analysis['file_name']}")
        summary_parts.append(f"• Size: {analysis['shape'][0]:,} rows × {analysis['shape'][1]} columns")
        summary_parts.append(f"• Data Quality: {analysis['data_quality']['complete_rows']:,} complete rows ({(analysis['data_quality']['complete_rows']/analysis['shape'][0]*100):.1f}%)")
        
        if analysis['data_quality']['duplicate_rows'] > 0:
            summary_parts.append(f"• Duplicates: {analysis['data_quality']['duplicate_rows']} rows")
        
        # Column analysis
        summary_parts.append(f"\n🔢 COLUMN ANALYSIS:")
        summary_parts.append(f"• Numeric columns ({len(analysis['numeric_columns'])}): {', '.join(analysis['numeric_columns'][:5])}{'...' if len(analysis['numeric_columns']) > 5 else ''}")
        summary_parts.append(f"• Categorical columns ({len(analysis['categorical_columns'])}): {', '.join(analysis['categorical_columns'][:5])}{'...' if len(analysis['categorical_columns']) > 5 else ''}")
        
        if analysis['datetime_columns']:
            summary_parts.append(f"• DateTime columns: {', '.join(analysis['datetime_columns'])}")
        
        # Statistical insights for numeric data
        if analysis['numeric_insights']:
            summary_parts.append(f"\n📈 STATISTICAL INSIGHTS:")
            for col, stats in list(analysis['numeric_insights'].items())[:3]:  # Limit to 3 columns
                summary_parts.append(
                    f"• {col}: mean={stats['mean']:.2f}, range=[{stats['min']:.2f}, {stats['max']:.2f}], "
                    f"std={stats['std']:.2f}, outliers={stats['outliers_count']}"
                )
                if abs(stats['skewness']) > 1:
                    summary_parts.append(f"  └ Data is {'positively' if stats['skewness'] > 0 else 'negatively'} skewed ({stats['skewness']:.2f})")
        
        # Categorical insights
        if analysis['categorical_insights']:
            summary_parts.append(f"\n📝 CATEGORICAL INSIGHTS:")
            for col, stats in list(analysis['categorical_insights'].items())[:3]:  # Limit to 3 columns
                summary_parts.append(
                    f"• {col}: {stats['unique_count']} unique values, "
                    f"most frequent='{stats['most_frequent']}' ({stats['concentration']:.1f}%)"
                )
        
        # Correlation insights
        if analysis['strong_correlations']:
            summary_parts.append(f"\n🔗 CORRELATION INSIGHTS:")
            for corr in analysis['strong_correlations'][:3]:  # Show top 3 correlations
                summary_parts.append(
                    f"• {corr['var1']} ↔ {corr['var2']}: {corr['correlation']} ({corr['strength']})"
                )
        
        # Data patterns
        if analysis['data_patterns']:
            patterns = analysis['data_patterns']
            summary_parts.append(f"\n🎯 DATA PATTERNS:")
            
            if patterns['potential_time_series']:
                summary_parts.append(f"• Time-based columns detected: {', '.join(patterns['potential_time_series'])}")
            
            if patterns['data_skewness']:
                skewed_cols = list(patterns['data_skewness'].keys())[:2]
                summary_parts.append(f"• Skewed distributions: {', '.join(skewed_cols)}")
            
            if patterns['suggested_groupings']:
                top_grouping = patterns['suggested_groupings'][0]
                summary_parts.append(f"• Recommended analysis: {top_grouping['group_by']} vs {top_grouping['analyze']}")
        
        # Missing data analysis
        missing_cols = [col for col, count in analysis['missing_values'].items() if count > 0]
        if missing_cols:
            summary_parts.append(f"\n⚠️ MISSING DATA:")
            for col in missing_cols[:3]:  # Show top 3 with missing data
                count = analysis['missing_values'][col]
                percentage = analysis['data_quality']['missing_percentage_by_column'][col]
                summary_parts.append(f"• {col}: {count} missing ({percentage:.1f}%)")
        
        # Sample data for context
        summary_parts.append(f"\n📋 SAMPLE DATA (first 2 rows):")
        for i, row in enumerate(analysis['sample_data'][:2]):
            row_str = ', '.join([f"{k}={v}" for k, v in list(row.items())[:4]])  # Show first 4 columns
            summary_parts.append(f"• Row {i+1}: {row_str}{'...' if len(row) > 4 else ''}")
        
        return "\n".join(summary_parts)

    def generate_insights_with_ai(self, analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Generate insights and presentation structure using AI with enhanced context"""
        
        # Build comprehensive data summary for AI
        data_summary = self._build_comprehensive_data_summary(analysis)
        print("\n Build comprehensive data summary for AI:\n", data_summary)
        prompt = (
            "You are an expert data analyst creating a comprehensive presentation from this CSV dataset analysis.\n\n"
            f"{data_summary}\n\n"
            "ANALYSIS REQUIREMENTS:\n"
            "• Focus on the most significant patterns, correlations, and insights\n"
            "• Consider data quality issues (missing values, outliers, skewness)\n"
            "• Identify business-relevant findings and actionable insights\n"
            "• Use statistical evidence to support your insights\n"
            "• Keep insights concise and slide-friendly (max 100 characters per insight)\n"
            "• Include multiple points in insights for depth and clarity\n"
            "• Consider the context and potential use cases for this data\n"
            "• Generate comprehensive overview slides with 5-8 bullet points each\n"
            "• Include detailed findings, trends, and recommendations\n\n"
            "CHART SELECTION STRATEGY:\n"
            "• Bar charts: For comparing categories or showing distributions\n"
            "• Pie charts: For proportional analysis (limit to 5-8 categories)\n"
            "• Line charts: For trends, time series, or ordered data\n"
            "• Scatter plots: For exploring relationships between numeric variables\n"
            "• Heatmaps: For correlation matrices or multi-dimensional comparisons\n"
            "• Select charts that best reveal the data's story and patterns\n\n"
            "OUTPUT FORMAT - Return _only_ this JSON structure:\n"
            "{\n"
            '  "title": "Compelling, concise presentation title (max 60 chars)",\n'
            '  "insights": [\n'
            '    "Key insight with data evidence (max 100 chars)",\n'
            '    "Business finding with statistical support (max 100 chars)",\n'
            '    "Pattern discovery with implications (max 100 chars)",\n'
            '    "Data quality observation (max 100 chars)",\n'
            '    "Trend analysis with business impact (max 100 chars)",\n'
            '    "Statistical finding with actionable recommendation (max 100 chars)"\n'
            '  ],\n'
            '  "slides": [\n'
            '    {\n'
            '      "title": "Executive Summary (max 50 chars)",\n'
            '      "content": [\n'
            '        "Dataset contains X records with Y key metrics analyzed",\n'
            '        "Primary finding: [most significant pattern/correlation]",\n'
            '        "Key opportunity: [business recommendation]",\n'
            '        "Data quality: [completeness percentage and main issues]",\n'
            '        "Geographic/temporal focus: [main patterns]",\n'
            '        "Critical insight: [most important business implication]",\n'
            '        "Recommended action: [specific next steps]"\n'
            '      ],\n'
            '      "slide_type": "overview"\n'
            '    },\n'
            '    {\n'
            '      "title": "Key Findings & Insights (max 50 chars)",\n'
            '      "content": [\n'
            '        "Statistical correlation: [strongest relationship found]",\n'
            '        "Distribution pattern: [skewness/outlier insights]",\n'
            '        "Categorical analysis: [dominant categories/frequencies]",\n'
            '        "Trend analysis: [temporal patterns if applicable]",\n'
            '        "Business impact: [financial/operational implications]",\n'
            '        "Risk factors: [potential issues identified]",\n'
            '        "Growth opportunities: [areas for improvement]"\n'
            '      ],\n'
            '      "slide_type": "insights"\n'
            '    },\n'
            '    {\n'
            '      "title": "Data Quality Assessment (max 50 chars)",\n'
            '      "content": [\n'
            '        "Completeness: [percentage] complete with [missing count] missing values",\n'
            '        "Outliers: [count] extreme values detected and handled",\n'
            '        "Data types: [numeric count] numeric, [text count] categorical columns",\n'
            '        "Duplicates: [count] duplicate records identified and removed",\n'
            '        "Consistency: [validation results] business rule compliance",\n'
            '        "Reliability: Data quality score of [percentage] achieved"\n'
            '      ],\n'
            '      "slide_type": "quality"\n'
            '    }\n'
            '  ],\n'
            '  "recommended_charts": [\n'
            '    {\n'
            '      "type": "bar|pie|line|scatter|heatmap",\n'
            '      "x_column": "specific_column_name",\n'
            '      "y_column": "specific_column_name_or_null",\n'
            '      "title": "Specific, descriptive chart title",\n'
            '      "purpose": "Why this chart reveals important insights"\n'
            '    }\n'
            '  ]\n'
            "}\n\n"
            "CRITICAL REQUIREMENTS:\n"
            "1. Generate 4-6 different chart types for comprehensive analysis\n"
            "2. Base insights on actual statistical evidence from the data\n"
            "3. Select columns strategically based on data characteristics\n"
            "4. Focus on charts that reveal the most important patterns\n"
            "5. Include both individual variable analysis and relationship exploration\n"
            "6. Ensure each chart serves a specific analytical purpose\n"
            "7. Prioritize actionable business insights over basic descriptions"
        )
        try:
            resp = self.openai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a data-analyst assistant. Output ONLY valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.2,
                max_tokens=1500
            )
            raw = resp.choices[0].message.content
            print("🔍 RAW AI RESPONSE:\n", raw)
            match = re.search(r'\{.*\}', raw, re.DOTALL)
            if not match:
                return self._get_fallback_structure(analysis)
            
            ai_result = json.loads(match.group(0))
            
            # Ensure we have multiple chart recommendations
            if 'recommended_charts' not in ai_result or len(ai_result['recommended_charts']) < 3:
                print("🔧 AI provided insufficient chart recommendations. Adding smart defaults.")
                smart_charts = self._get_smart_chart_recommendations(analysis)
                ai_result['recommended_charts'] = smart_charts
            
            # Also ensure we have enough insights
            if 'insights' not in ai_result or len(ai_result['insights']) < 3:
                ai_result['insights'] = [
                    f"Dataset contains {analysis['shape'][0]:,} records across {analysis['shape'][1]} columns",
                    f"Found {len(analysis['numeric_columns'])} numeric and {len(analysis['categorical_columns'])} categorical variables",
                    f"Generated {len(ai_result['recommended_charts'])} comprehensive visualizations",
                    "Analysis reveals key patterns and relationships in the data"
                ]
            
            return ai_result
        except Exception as e:
            print(f"Error getting AI insights: {e}")
            return self._get_fallback_structure(analysis)

    def _get_smart_chart_recommendations(self, analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Generate smart chart recommendations based on data characteristics"""
        nums = analysis['numeric_columns']
        cats = analysis['categorical_columns']
        recommendations = []
        
        # 1. Bar chart - for distribution analysis
        if cats and nums:
            # Bar chart with categorical x and numeric y
            recommendations.append({
                "type": "bar",
                "x_column": cats[0],
                "y_column": nums[0],
                "title": f"{cats[0]} vs {nums[0]}",
                "purpose": "Show relationship between category and numeric data"
            })
        elif nums:
            # Bar chart for numeric distribution (histogram style)
            recommendations.append({
                "type": "bar",
                "x_column": nums[0],
                "y_column": None,
                "title": f"{nums[0]} Distribution",
                "purpose": "Show distribution of numeric data"
            })
        elif cats:
            # Bar chart for categorical distribution
            recommendations.append({
                "type": "bar",
                "x_column": cats[0],
                "y_column": None,
                "title": f"{cats[0]} Distribution",
                "purpose": "Show category distribution"
            })
        
        # 2. Pie chart - for categorical data proportions
        if cats and nums:
            # Pie chart with categorical data and numeric values
            recommendations.append({
                "type": "pie",
                "x_column": cats[0],
                "y_column": nums[0],
                "title": f"{cats[0]} by {nums[0]}",
                "purpose": "Show proportional breakdown with values"
            })
        elif cats:
            # Pie chart for categorical proportions
            recommendations.append({
                "type": "pie",
                "x_column": cats[0],
                "y_column": None,
                "title": f"{cats[0]} Proportions",
                "purpose": "Show proportional breakdown"
            })
        
        # 3. Line chart - for trend analysis
        if len(nums) >= 1:
            recommendations.append({
                "type": "line",
                "x_column": nums[0] if len(nums) >= 1 else None,
                "y_column": nums[1] if len(nums) >= 2 else nums[0],
                "title": "Trend Analysis",
                "purpose": "Show data trends over sequence"
            })
        
        # 4. Scatter plot - for correlation between numeric columns
        if len(nums) >= 2:
            recommendations.append({
                "type": "scatter",
                "x_column": nums[0],
                "y_column": nums[1],
                "title": f"{nums[0]} vs {nums[1]}",
                "purpose": "Show correlation between variables"
            })
        
        # 5. Heatmap - for correlation matrix
        if len(nums) >= 2:
            recommendations.append({
                "type": "heatmap",
                "x_column": None,
                "y_column": None,
                "title": "Correlation Matrix",
                "purpose": "Show correlations between all numeric variables"
            })
        
        # 6. Additional mixed analysis if we have both types
        if cats and nums and len(cats) > 1:
            recommendations.append({
                "type": "bar",
                "x_column": cats[1] if len(cats) > 1 else cats[0],
                "y_column": nums[1] if len(nums) > 1 else nums[0],
                "title": f"{cats[1] if len(cats) > 1 else cats[0]} vs {nums[1] if len(nums) > 1 else nums[0]}",
                "purpose": "Show additional relationship between category and numeric data"
            })
        
        return recommendations[:5]  # Limit to 5 charts

    def _get_fallback_structure(self, analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Fallback presentation structure with multiple chart types"""
        nums = analysis['numeric_columns']
        cats = analysis['categorical_columns']
        
        # Get smart chart recommendations
        chart_recommendations = self._get_smart_chart_recommendations(analysis)
        
        return {
            "title": f"Data Analysis Report: {analysis['file_name']}",
            "insights": [
                f"Dataset contains {analysis['shape'][0]:,} records across {analysis['shape'][1]} columns",
                f"Found {len(nums)} numeric and {len(cats)} categorical variables",
                f"Generated {len(chart_recommendations)} different visualizations",
                "Comprehensive analysis reveals key data patterns and distributions"
            ],
            "slides": [
                {"title": "Title Slide", "content": [], "slide_type": "title"},
                {"title": "Executive Summary", "content": [
                    f"📊 Dataset contains {analysis['shape'][0]:,} records across {analysis['shape'][1]} key variables",
                    f"🔢 Analysis covers {len(nums)} numeric and {len(cats)} categorical dimensions",
                    f"📈 Generated {len(chart_recommendations)} comprehensive visualizations for insights",
                    f"📋 Data quality achieved: {((analysis['shape'][0] * analysis['shape'][1] - sum(analysis['missing_values'].values())) / (analysis['shape'][0] * analysis['shape'][1]) * 100):.1f}% completeness",
                    f"🎯 Key focus areas: {', '.join(nums[:3])} for quantitative analysis",
                    f"📝 Primary categories: {', '.join(cats[:3])} for segmentation",
                    f"💡 Ready for comprehensive business intelligence and decision support"
                ], "slide_type": "overview"},
                {"title": "Key Findings & Insights", "content": [
                    f"📊 Statistical analysis reveals {len(analysis.get('strong_correlations', []))} significant correlations",
                    f"🔍 Data distribution shows patterns across {len(analysis['categorical_insights'])} key categories",
                    f"📈 Quantitative metrics span from {min([analysis['numeric_insights'][col]['min'] for col in nums[:3]]) if nums else 'N/A'} to {max([analysis['numeric_insights'][col]['max'] for col in nums[:3]]) if nums else 'N/A'}",
                    f"⚠️ Data quality considerations: {sum(analysis['missing_values'].values())} missing values identified",
                    f"🎯 Business impact: Actionable insights across {len(analysis['categorical_columns'])} operational dimensions",
                    f"📋 Trend analysis: Time-based patterns {'detected' if analysis['datetime_columns'] else 'not available'}",
                    f"💼 Strategic recommendations: Focus on top-performing segments and outlier management"
                ], "slide_type": "insights"}
            ],
            "recommended_charts": chart_recommendations
        }

    
    def create_chart_from_data(self, chart_config: Dict[str, Any]) -> str:
        """Create chart from actual CSV data with proper axis formatting"""
        plt.figure(figsize=(12, 8))
        plt.clf()
        
        chart_type = chart_config.get('chart_type', 'bar')
        title = chart_config.get('title', 'Data Chart')
        
        try:
            if chart_type == 'bar':
                self._create_bar_chart(chart_config)
            elif chart_type == 'pie':
                self._create_pie_chart(chart_config)
            elif chart_type == 'line':
                self._create_line_chart(chart_config)
            elif chart_type == 'scatter':
                self._create_scatter_chart(chart_config)
            elif chart_type == 'heatmap':
                self._create_heatmap_chart(chart_config)
            else:
                self._create_default_chart(chart_config)
            
            # Apply common formatting for non-pie charts
            if chart_type != 'pie' and chart_type != 'heatmap':
                self._format_chart_axes(chart_config)
            
            plt.title(title, fontsize=16, fontweight='bold', pad=20)
            plt.tight_layout(pad=2.0)  # More padding for better fit
            
            # Save chart
            chart_filename = f"chart_{len(self.charts_created)}_{chart_type}.png"
            chart_path = os.path.join(os.path.dirname(__file__), chart_filename)
            plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
            plt.close()
            
            self.charts_created.append(chart_path)
            return chart_path
            
        except Exception as e:
            print(f"Error creating chart: {e}")
            return self._create_fallback_chart(title)
    
    def _create_bar_chart(self, config: Dict[str, Any]):
        """Create bar chart from data"""
        x_col = config.get('x_column')
        y_col = config.get('y_column')
        
        if x_col and y_col and x_col in self.df.columns and y_col in self.df.columns:
            # Group and aggregate data if needed
            if self.df[x_col].dtype == 'object':
                data = self.df.groupby(x_col)[y_col].sum().head(10)
                plt.bar(range(len(data)), data.values, color=plt.cm.Set3(range(len(data))))
                plt.xticks(range(len(data)), data.index, rotation=45)
                plt.ylabel(y_col)
            else:
                plt.hist(self.df[x_col].dropna(), bins=20, color='skyblue', alpha=0.7)
                plt.xlabel(x_col)
                plt.ylabel('Frequency')
        else:
            # Create chart with numeric columns
            numeric_cols = self.df.select_dtypes(include=[np.number]).columns[:5]
            if len(numeric_cols) > 0:
                means = [self.df[col].mean() for col in numeric_cols]
                plt.bar(numeric_cols, means, color=plt.cm.Set3(range(len(numeric_cols))))
                plt.xticks(rotation=45)
                plt.ylabel('Average Values')
    
    def _create_pie_chart(self, config: Dict[str, Any]):
        """Create pie chart from data"""
        x_col = config.get('x_column')
        
        if x_col and x_col in self.df.columns:
            if self.df[x_col].dtype == 'object':
                value_counts = self.df[x_col].value_counts().head(8)
                plt.pie(value_counts.values, labels=value_counts.index, autopct='%1.1f%%', startangle=90)
            else:
                # Create bins for numeric data
                bins = pd.cut(self.df[x_col].dropna(), bins=5)
                value_counts = bins.value_counts()
                plt.pie(value_counts.values, labels=[str(x) for x in value_counts.index], autopct='%1.1f%%')
        else:
            # Default: show data types distribution
            dtypes_count = self.df.dtypes.value_counts()
            plt.pie(dtypes_count.values, labels=dtypes_count.index, autopct='%1.1f%%', startangle=90)
    
    def _create_line_chart(self, config: Dict[str, Any]):
        """Create line chart from data"""
        x_col = config.get('x_column')
        y_col = config.get('y_column')
        
        if x_col and y_col and x_col in self.df.columns and y_col in self.df.columns:
            # Sort by x column
            sorted_data = self.df[[x_col, y_col]].dropna().sort_values(x_col)
            plt.plot(sorted_data[x_col], sorted_data[y_col], marker='o', linewidth=2, markersize=6)
            plt.xlabel(x_col)
            plt.ylabel(y_col)
            plt.grid(True, alpha=0.3)
        else:
            # Default: show trend of first numeric column
            numeric_cols = self.df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                plt.plot(self.df.index, self.df[col], marker='o', linewidth=2)
                plt.xlabel('Index')
                plt.ylabel(col)
                plt.grid(True, alpha=0.3)
    
    def _create_scatter_chart(self, config: Dict[str, Any]):
        """Create scatter plot from data"""
        x_col = config.get('x_column')
        y_col = config.get('y_column')
        
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns
        
        if x_col and y_col and x_col in numeric_cols and y_col in numeric_cols:
            plt.scatter(self.df[x_col], self.df[y_col], alpha=0.6, s=50)
            plt.xlabel(x_col)
            plt.ylabel(y_col)
        elif len(numeric_cols) >= 2:
            plt.scatter(self.df[numeric_cols[0]], self.df[numeric_cols[1]], alpha=0.6, s=50)
            plt.xlabel(numeric_cols[0])
            plt.ylabel(numeric_cols[1])
    
    def _create_heatmap_chart(self, config: Dict[str, Any]):
        """Create heatmap from data"""
        numeric_data = self.df.select_dtypes(include=[np.number])
        
        if len(numeric_data.columns) > 1:
            correlation_matrix = numeric_data.corr()
            sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', center=0, square=True)
        else:
            # Create a simple heatmap with data summary
            summary_data = self.df.describe().T
            if not summary_data.empty:
                sns.heatmap(summary_data[['mean', 'std']], annot=True, cmap='viridis')
    
    def _format_chart_axes(self, config: Dict[str, Any]):
        """Format chart axes with proper labels, units, and scaling"""
        x_col = config.get('x_column')
        y_col = config.get('y_column')
        
        # Get current axes
        ax = plt.gca()
        
        # Format X-axis
        if x_col:
            # Add proper X-axis label with units if numeric
            if x_col in self.df.select_dtypes(include=[np.number]).columns:
                x_unit = self._detect_unit(x_col, self.df[x_col])
                plt.xlabel(f"{x_col}{x_unit}", fontsize=12, fontweight='bold')
                
                # Format X-axis ticks for better readability
                x_values = self.df[x_col].dropna()
                if x_values.max() > 1000:
                    ax.xaxis.set_major_formatter(plt.FuncFormatter(self._format_large_numbers))
            else:
                plt.xlabel(x_col, fontsize=12, fontweight='bold')
                # Rotate labels if they're long text
                plt.xticks(rotation=45, ha='right')
        
        # Format Y-axis
        if y_col:
            # Add proper Y-axis label with units if numeric
            if y_col in self.df.select_dtypes(include=[np.number]).columns:
                y_unit = self._detect_unit(y_col, self.df[y_col])
                plt.ylabel(f"{y_col}{y_unit}", fontsize=12, fontweight='bold')
                
                # Format Y-axis ticks for better readability
                y_values = self.df[y_col].dropna()
                if y_values.max() > 1000:
                    ax.yaxis.set_major_formatter(plt.FuncFormatter(self._format_large_numbers))
            else:
                plt.ylabel(y_col, fontsize=12, fontweight='bold')
        elif not y_col and 'frequency' in str(ax.get_ylabel()).lower():
            # For histograms and frequency charts
            plt.ylabel('Frequency (Count)', fontsize=12, fontweight='bold')
        
        # Add grid for better readability
        ax.grid(True, alpha=0.3, linestyle='--')
        
        # Ensure proper margins and spacing
        plt.subplots_adjust(bottom=0.15, left=0.15, right=0.95, top=0.9)
        
        # Format tick labels for better readability
        ax.tick_params(axis='both', which='major', labelsize=10)
        
        # Set axis limits if needed to prevent overcrowding
        if x_col and x_col in self.df.select_dtypes(include=[np.number]).columns:
            x_data = self.df[x_col].dropna()
            x_range = x_data.max() - x_data.min()
            if x_range > 0:
                plt.xlim(x_data.min() - x_range * 0.05, x_data.max() + x_range * 0.05)
        
        if y_col and y_col in self.df.select_dtypes(include=[np.number]).columns:
            y_data = self.df[y_col].dropna()
            y_range = y_data.max() - y_data.min()
            if y_range > 0:
                plt.ylim(y_data.min() - y_range * 0.05, y_data.max() + y_range * 0.05)
    
    def _detect_unit(self, column_name: str, data: pd.Series) -> str:
        """Detect appropriate unit indicator for numeric data"""
        col_lower = column_name.lower()
        max_val = data.max()
        
        # Common unit patterns
        if any(word in col_lower for word in ['price', 'cost', 'revenue', 'profit', 'sales', 'amount']):
            if max_val >= 1000000:
                return " ($ Millions)"
            elif max_val >= 1000:
                return " ($ Thousands)"
            else:
                return " ($)"
        
        elif any(word in col_lower for word in ['percent', 'rate', '%']):
            return " (%)"
        
        elif any(word in col_lower for word in ['count', 'number', 'qty', 'quantity']):
            if max_val >= 1000000:
                return " (Millions)"
            elif max_val >= 1000:
                return " (Thousands)"
            else:
                return " (Count)"
        
        elif any(word in col_lower for word in ['time', 'duration', 'hours', 'minutes']):
            return " (Time)"
        
        elif any(word in col_lower for word in ['weight', 'mass']):
            return " (kg)"
        
        elif any(word in col_lower for word in ['distance', 'length', 'height']):
            return " (m)"
        
        elif any(word in col_lower for word in ['temperature', 'temp']):
            return " (°C)"
        
        else:
            # Generic unit based on magnitude
            if max_val >= 1000000:
                return " (Millions)"
            elif max_val >= 1000:
                return " (Thousands)"
            else:
                return ""
    
    def _format_large_numbers(self, x, pos):
        """Format large numbers for axis labels"""
        if x >= 1000000:
            return f'{x/1000000:.1f}M'
        elif x >= 1000:
            return f'{x/1000:.1f}K'
        else:
            return f'{x:.0f}'
    
    def _create_default_chart(self, config: Dict[str, Any]):
        """Create default chart when specific type fails"""
        numeric_cols = self.df.select_dtypes(include=[np.number]).columns[:5]
        if len(numeric_cols) > 0:
            means = [self.df[col].mean() for col in numeric_cols]
            plt.bar(numeric_cols, means, color='skyblue')
            plt.xticks(rotation=45)
            plt.ylabel('Average Values')
        else:
            plt.text(0.5, 0.5, 'No suitable data for visualization', 
                    horizontalalignment='center', verticalalignment='center', 
                    transform=plt.gca().transAxes, fontsize=14)
    
    def _create_fallback_chart(self, title: str) -> str:
        """Create a fallback chart when errors occur"""
        plt.figure(figsize=(10, 6))
        plt.text(0.5, 0.5, f'Chart: {title}\n(Error in data processing)', 
                horizontalalignment='center', verticalalignment='center', 
                fontsize=14)
        
        chart_filename = f"fallback_chart_{len(self.charts_created)}.png"
        chart_path = os.path.join(os.path.dirname(__file__), chart_filename)
        plt.savefig(chart_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        self.charts_created.append(chart_path)
        return chart_path

    def create_presentation_from_csv(self, file_path: str, output_filename: str = None, sheet_name: str = None, named_range: str = None) -> str:
        """Complete workflow: analyze CSV/Excel and create presentation"""
        file_type = self.detect_file_type(file_path)
        print(f"📊 Loading and analyzing {file_type.upper()} file: {file_path}")
        
        # For Excel files, show helpful information
        if file_type == 'excel' and not sheet_name:
            try:
                excel_info = self.load_excel_info(file_path)
                print(f"📋 Excel file contains {excel_info['total_sheets']} sheets:")
                for name, info in excel_info['sheets'].items():
                    status = "✅" if info['has_data'] else "❌"
                    print(f"  {status} {name}: {info['estimated_records']} rows")
            except Exception as e:
                print(f"⚠️  Could not read Excel info: {e}")
        
        # Load and analyze data (supports both CSV and Excel)
        analysis = self.load_and_analyze_data(file_path, sheet_name, named_range)
        

        print(f"🤖 Generating insights with AI...")
        structure = self.generate_insights_with_ai(analysis)

        if output_filename is None:
            base = os.path.splitext(analysis["file_name"])[0]
            output_filename = f"{base}_analysis_presentation.pptx"

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 1. Core slides (title, overview, chart, insights, etc.)
        for slide in structure.get("slides", []):
            stype = slide.get("slide_type", "content")
            if stype == "title":
                self._create_title_slide(prs, slide, structure)
            elif stype == "chart":
                self._create_chart_slide(prs, slide)
            else:
                self._create_content_slide(prs, slide)

        # 2. If AI succeeded, structure["recommended_charts"] holds multiple specs.
        #    Generate one slide per recommended chart (bar, pie, line, scatter, heatmap, …)
        for rec in structure.get("recommended_charts", []):
            slide_data = {
                "title": rec.get("title", ""),
                "slide_type": "chart",
                "chart_config": {
                    "chart_type": rec.get("type", ""),
                    "x_column": rec.get("x_column"),
                    "y_column": rec.get("y_column"),
                    "title": rec.get("title", "")
                }
            }
            self._create_chart_slide(prs, slide_data)

        prs.save(output_filename)
        self._cleanup_chart_files()
        print(f"✅ Presentation saved as: {output_filename}")
        return output_filename

    # ... [the rest of your helper methods: _create_title_slide,
    #      _create_chart_slide, _create_content_slide, _cleanup_chart_files] ...
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict[str, Any], structure: Dict[str, Any]):
        """Create title slide with proper positioning to avoid overlaps"""
        slide_layout = prs.slide_layouts[6]  # Use blank layout for full control
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title with manual positioning
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
        title_frame = title_shape.text_frame
        title_frame.text = slide_data['title']
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        
        # Style title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(40)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(44, 62, 80)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add subtitle with proper spacing
        subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(2))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.margin_left = Inches(0)
        subtitle_frame.margin_right = Inches(0)
        subtitle_frame.margin_top = Inches(0)
        subtitle_frame.margin_bottom = Inches(0)
        
        subtitle_text = "\n".join([
            f"📊 Dataset: {self.data_analysis['shape'][0]:,} rows, {self.data_analysis['shape'][1]} columns",
            f"🎯 Key Insights: {len(structure.get('insights', []))} findings",
            f"📅 Generated: {datetime.now().strftime('%B %d, %Y')}"
        ])
        subtitle_frame.text = subtitle_text
        
        # Style subtitle
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(127, 140, 141)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_after = Pt(6)
    
    def _create_chart_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Create slide with chart using blank layout to avoid overlaps"""
        slide_layout = prs.slide_layouts[6]  # Use blank layout for full control
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title with proper positioning and margins
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = slide_data['title']
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        
        # Style title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(44, 62, 80)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Create and add chart with proper spacing from title
        chart_config = slide_data.get('chart_config', {})
        chart_path = self.create_chart_from_data(chart_config)
        
        # Add chart image to slide with proper positioning
        slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.5))
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Create content slide with controlled content length to prevent overflow"""
        slide_layout = prs.slide_layouts[6]  # Use blank layout for consistency
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title with manual positioning
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = self._truncate_text(slide_data['title'], 80)  # Limit title length
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        title_frame.word_wrap = True
        
        # Style title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(26)  # Slightly smaller to ensure fit
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(44, 62, 80)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add content with proper spacing from title and controlled height
        content_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.73), Inches(5.4))
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.2)
        content_frame.margin_right = Inches(0.2)
        content_frame.margin_top = Inches(0.2)
        content_frame.margin_bottom = Inches(0.2)
        content_frame.word_wrap = True
        content_frame.auto_size = None  # Disable auto-sizing to prevent overflow
        
        # Limit content to prevent overflow (max 10 bullet points)
        content_items = slide_data.get('content', [])
        max_bullets = 10
        limited_content = content_items[:max_bullets]
        
        # Add bullet points with controlled length
        for i, bullet_point in enumerate(limited_content):
            # Truncate long bullet points to prevent text overflow
            truncated_bullet = self._truncate_text(str(bullet_point), 120)
            
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {truncated_bullet}"
            p.font.size = Pt(14)  # Slightly smaller for better fit
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.space_after = Pt(8)  # Reduced spacing
            p.space_before = Pt(4)
            p.line_spacing = 1.1  # Tighter line spacing
        
        # Add "...and more" indicator if content was truncated
        if len(content_items) > max_bullets:
            p = content_frame.add_paragraph()
            p.text = f"• ...and {len(content_items) - max_bullets} more insights"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(127, 140, 141)  # Lighter color
            p.font.italic = True
    
    def _truncate_text(self, text: str, max_length: int) -> str:
        """Truncate text to prevent overflow with proper word boundaries"""
        if len(text) <= max_length:
            return text
        
        # Find the last space before max_length to avoid cutting words
        truncated = text[:max_length]
        last_space = truncated.rfind(' ')
        
        if last_space > max_length * 0.7:  # If space is reasonably close to end
            return truncated[:last_space] + "..."
        else:
            return truncated + "..."
    
    def _cleanup_chart_files(self):
        """Clean up temporary chart files"""
        for chart_path in self.charts_created:
            try:
                if os.path.exists(chart_path):
                    os.remove(chart_path)
            except Exception as e:
                print(f"Warning: Could not remove chart file {chart_path}: {e}")
        
        self.charts_created = []

def main():
    import argparse
    parser = argparse.ArgumentParser(description="Generate PPT from CSV or Excel files")
    parser.add_argument('file', help="Path to the CSV or Excel file")
    parser.add_argument('-o', '--output', help="Output .pptx filename")
    parser.add_argument('-s', '--sheet', help="Excel sheet name (if not specified, auto-selects best sheet)")
    parser.add_argument('-r', '--range', help="Named range in Excel file (optional)")
    parser.add_argument('--list-sheets', action='store_true', help="List all sheets in Excel file and exit")
    args = parser.parse_args()

    try:
        gen = CSVPPTGenerator()
        
        # Special case: just list sheets and exit
        if args.list_sheets:
            file_type = gen.detect_file_type(args.file)
            if file_type == 'excel':
                excel_info = gen.load_excel_info(args.file)
                print(f"\n📁 Excel file: {args.file}")
                print(f"📊 Total sheets: {excel_info['total_sheets']}")
                print(f"✅ Sheets with data: {excel_info['sheets_with_data']}")
                print("\n📋 Sheet Details:")
                for sheet_name, info in excel_info['sheets'].items():
                    status = "✅ Has data" if info['has_data'] else "❌ Empty"
                    print(f"  • {sheet_name}: {info['estimated_records']} rows - {status}")
                
                if excel_info['named_ranges']:
                    print("\n🎯 Named Ranges:")
                    for nr in excel_info['named_ranges']:
                        print(f"  • {nr['name']}: {nr['range']} (Sheet: {nr['sheet']})")
                else:
                    print("\n🎯 No named ranges found")
            else:
                print(f"❌ File is not Excel format: {args.file}")
            return
        
        # Generate presentation
        gen.create_presentation_from_csv(args.file, args.output, args.sheet, args.range)
        
    except Exception as e:
        print(f"❌ Error: {e}")
        exit(1)

if __name__ == "__main__":
    main()
